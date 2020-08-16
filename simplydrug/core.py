# AUTOGENERATED! DO NOT EDIT! File to edit: simplydrug.ipynb (unless otherwise specified).

__all__ = ['exception_handler', 'add_layout', 'order_wells', 'hts_heatmap', 'sum_statistics', 'normalize_z',
           'histogram_feature', 'plot_treatments', 'plot_curve_raw', 'plot_curve_mean', 'pointplot_plate',
           'calculate_growth_score', 'filter_curves', 'inv_log', 'll4', 'pDose', 'dose_response', 'dose_response_old',
           'plot_dr_viability', 'prune_compound', 'plot_polynomial', 'df_to_table', 'create_presentation']

# Cell
def exception_handler(func):
    """Exception handler helper function.
    :param func: any function."""

    def wrapper_func(*args, **kwargs):
        try:
            return func(*args, **kwargs)
        except Exception as e:
            print(f'{func.__name__} exception: ')
            print(e)
            return None
    return wrapper_func

# Cell
@exception_handler
def add_layout(df = None, layout_path = None, chem_path = None, chem_plate = None):
    """Takes DataFrame, the path to layout excel file, chemical library path, and plate name and return
    DataFrame with added layout columns. In the layout file, each sheet should represent different aspects
    of the layout. Obligatory sheets are 'Well' (for well names), and 'Status'. 'Status' column defines
    each well as 'Sample',  'Positive' or 'Negative' control, or 'Reference'.
    'Reference' wells are usually excluded from the analysis.
    Sheet names are translated to column names in the final DataFrame.
    :param df: pandas DataFrame with the data.
    :param layout_path: path to layout xlsx file.
    :param chem_path: path to chemical library file, optional.
    :param chem_library: name of the chemical library plate, optional.
    :return: DataFrame with added layout columns.
    """

    import pandas as pd
    import numpy as np

    layout = pd.DataFrame()
    for sheet in pd.ExcelFile(layout_path).sheet_names:   # create columns from excel file
        layout[sheet] = np.asarray(pd.ExcelFile(layout_path).parse(sheet)).reshape(-1)
        print('Added ', sheet)

    if chem_path and chem_plate: # add compound IDs
        compounds = pd.read_csv(chem_path, low_memory = False)
        layout = pd.merge(layout, compounds[compounds.Plate == chem_plate], how = 'left', on = 'Well')
        print('Added compounds: ',chem_plate,'\n')

    else:
        print('Chemical library not requested')

    output = pd.merge(df, layout, how = 'left', on = 'Well')
    return output

# Cell
@exception_handler
def order_wells(x = None):
    import re
    """Orders wells as they appear in the plate.
    :param x: list of wells.
    :return: Ordered list of wells."""
    convert = lambda text: int(text) if text.isdigit() else text
    alphanum_key = lambda key: [ convert(c) for c in re.split('([0-9]+)', key) ]
    return sorted(x, key = alphanum_key)

# Cell
@exception_handler
def hts_heatmap(df = None, layout_path = None, features = None, save_as = None, path = None):
    """Takes DataFrame, list of features, a path to layout file and to the output folder and creates
    plate heatmap for input features.
    :param df: pandas DataFrame with the data.
    :param layout_path: path to layout xlsx file.
    :param features: list of features to build heatmaps on.
    :param save_as: filename to save the resulting figure, optional.
    :param path: path to the output folder to save the resulting figure, optional.
    :return: None. """

    import pandas as pd
    import numpy as np
    import matplotlib as mpl
    import matplotlib.colors as colors
    import matplotlib.pyplot as plt
    import seaborn as sns

    # merge data with well names and status
    layout = pd.DataFrame({'Well':np.asarray(pd.ExcelFile(layout_path).parse('Well')).reshape(-1),
                      'Status': np.asarray(pd.ExcelFile(layout_path).parse('Status')).reshape(-1)})
    data = pd.merge(layout, df.groupby('Well').mean(), how = 'left', on ='Well').set_index('Well', drop = False)
    samples = data[data['Status'] == 'Sample'].dropna()

    # define plate format
    if data.shape[0] == 384:
        plate_reshape = (16,24) # number of rows, number of columns for 384
        yticks=['A','B','C','D','E','F','G','H','I','J','K','L','M','N','O','P']
        xticks=['1','2','3','4','5','6','7','8','9','10','11','12','13','14','15','16','17','18','19','20','21','22', '23','24']
    elif data.shape[0] == 96:
        plate_reshape = (8,12) # number of rows, number of columns for 96
        yticks=['A','B','C','D','E','F','G','H']
        xticks=['1','2','3','4','5','6','7','8','9','10','11','12']
    else:
        print('Unknown plate format - cannot generate the heatmaps')

    # order wells properly
    ordered_data =  pd.DataFrame()
    for well in order_wells(data.Well):
        ordered_data = ordered_data.append(data.loc[well,:])

    # build heatmap for each feature
    for f in features:
        plate_view = ordered_data[f].values.reshape(plate_reshape)
        plate_view[np.isnan(plate_view)] = samples[f].mean() # fill missing values with sample mean value
        vmin = samples[f].mean() - 3*(samples[f].std()) # min value for the heatmap
        vmax = samples[f].mean() + 3*(samples[f].std()) # min value for the heatmap

        # plot
        ax = sns.heatmap(plate_view, vmin, vmax, center = samples[f].mean(), yticklabels = yticks, xticklabels = xticks, cmap = 'RdBu_r')
        ax.set_yticklabels(yticks, rotation = 0)
        ax.set_xticklabels(xticks,rotation = 0)
        ax.set_title(f + ' \n')
        fig = ax.get_figure()
        fig.set_size_inches(7, 4)
        if path:
            plt.savefig(path + '//' + f + save_as, bbox_inches = 'tight', dpi = 600)
        plt.show()
        plt.close()

# Cell
@exception_handler
def sum_statistics(df = None, feature = None):
    """Takes DataFrame and calculates summary statistics for the experiment. The data must
    contain the 'Status' column, defining each row as 'Sample', 'Positive' or 'Negative'
    control, or 'Reference'. 'Reference' wells are excluded from the analysis.
    :param df: pandas DataFrame with the data.
    :param feature: feature to calculate statistics.
    :return: summary statistics DataFrame.
    """
    import pandas as pd
    import numpy as np
    from scipy import stats

    st = None
    df = df[df.Status != 'Reference'][[feature, 'Status']]
    st = df.groupby(['Status']).agg([np.size, np.mean, np.std, np.var])
    st.columns = st.columns.droplevel()
    st['Feature'] = feature

    if 'Positive' in df['Status'].unique() and 'Negative' in df['Status'].unique():
        st['Z_factor'] = 1 - 3*(st.at['Positive','std'] + st.at['Negative','std'])/abs(st.at['Positive','mean'] - st.at['Negative','mean'])
        st['SB'] = st.at['Positive','mean']/st.at['Negative','mean']
        st = st.reset_index()[['Feature', 'Status', 'size', 'mean', 'std', 'var', 'Z_factor', 'SB' ]]
    else:
        print('sum_statistics: Failed calculate Z factor. Positive or Negative control is missing.')
        st = st.reset_index()[['Feature', 'Status', 'size', 'mean', 'std', 'var']]
        pass
    return st


# Cell
@exception_handler
def normalize_z(df = None, feature = None):
    """Takes DataFrame with experimental results and feature name and adds a column with
    normalized values of the feature.
    :param df: pandas DataFrame with the data.
    :param feature: name of the feature to normalize.
    :return: DataFrame with new column containing z-normalized feature values.
    """
    import pandas as pd
    mean = df[df['Status'] == 'Sample'][[feature]].mean()
    std = df[df['Status'] == 'Sample'][[feature]].std()
    df[feature + '_norm'] =  df[feature].apply(lambda x:(x - mean)/std)
    return(df)

# Cell
@exception_handler
def histogram_feature(df = None, feature = None, save_as = None, path = None):
    """Creates histogram of the input feature.
    :param df: pandas DataFrame with the data.
    :param feature: name of the feature to build histogram.
    :param save_as: filename to save the resulting figure, optional.
    :param path: path to the output folder to save the resulting figure, optional.
    :return: None."""
    import pandas as pd
    import matplotlib as mpl
    import matplotlib.colors as colors
    import matplotlib.pyplot as plt
    import seaborn as sns

    sns.distplot(df[feature].values)
    plt.plot([-2, -2], [0, 0.5], color = 'r', linestyle = '--', lw = 1.7)
    plt.plot([2, 2], [0, 0.5], color = 'r', linestyle = '--', lw = 1.7)
    if path:
        plt.savefig(path + '//' + save_as, bbox_inches = 'tight', dpi = 600)
    plt.show()
    plt.close()

# Cell
@exception_handler
def plot_treatments(df = None, x = None, y = None, column = None, kind = None, ylabel = None,
                    palette = None, height = None, aspect = None, save_as = None, path = None):
    """Creates plot by treatments. If your data has different treatments, set column = 'Treatment'.
    :param df: pandas DataFrame for plotting.
    :param x, y: names of variables in data, optional.
    :param column, kind, ylabel, palette, high, aspect: These features are seaborn.boxplot features.
    :param save_as: filename to save the resulting figure, optional.
    :param path: path to the output folder to save the resulting figure, optional.
    :return: None."""
    import pandas as pd
    import matplotlib as mpl
    import matplotlib.colors as colors
    import matplotlib.pyplot as plt
    import seaborn as sns

    plot_data = df[df['Status'] != 'Reference'] #filter out the Reference wells
    g = sns.catplot(x = x, y = y, data = plot_data, col = column, kind = kind, palette = palette,
                    height = height, aspect = aspect, margin_titles = False)
    axes = g.axes.flatten()
    axes[0].set_ylabel(ylabel)
    g.set_xticklabels(rotation = 90)
    if path:
        plt.savefig(path + '//' + save_as, bbox_inches = 'tight', dpi = 600)
    plt.show()
    plt.close()

# Cell
@exception_handler
def plot_curve_raw(df = None, x = None, y = None, units = None, hue = None, hue_order = None, xlabel = None,
                   ylabel = None, xlimit = None, palette = None, save_as = None, path = None):
    """Plot raw curves.
    :param df: pandas DataFrame for plotting.
    :param x, y: names of variables in data.
    :param units, hue, hue_order, xlabel, ylabel, xlimit, palette, high, aspect: These features are seaborn.boxplot features.
    :param save_as: filename to save the resulting figure, optional.
    :param path: path to the output folder to save the resulting figure, optional.
    :return: None."""
    import pandas as pd
    import matplotlib as mpl
    import matplotlib.colors as colors
    import matplotlib.pyplot as plt
    import seaborn as sns

    ax = sns.lineplot(data = df, x = x, y = y, units = units, hue = hue, hue_order = hue_order, palette = palette, estimator = None)
    ax.set_xlabel(xlabel)
    ax.set_ylabel(ylabel)
    ax.set_xlim(0, xlimit)
    fig = ax.get_figure()
    fig.set_size_inches(10, 7)

    if path:
        plt.savefig(path + '//' + save_as, bbox_inches = 'tight', dpi = 600)
    plt.show()
    plt.close()

# Cell
@exception_handler
def plot_curve_mean(df = None, x = None, y = None, hue = None, hue_order = None, xlabel = None, ylabel = None,
                    xlimit = None, palette = None, save_as = None, path = None):
    """Plot mean curves.
    :param df: pandas DataFrame for plotting.
    :param x, y: names of variables in data.
    :param hue, hue_order, xlabel, ylabel, xlimit, palette, high, aspect: These features are seaborn.boxplot features.
    :param save_as: filename to save the resulting figure, optional.
    :param path: path to the output folder to save the resulting figure, optional.
    :return: None.
    """
    import pandas as pd
    import matplotlib as mpl
    import matplotlib.colors as colors
    import matplotlib.pyplot as plt
    import seaborn as sns

    ax = sns.lineplot(data = df, x = x, y = y, hue = hue, hue_order = hue_order, palette = palette)
    sns.despine()
    ax.set_xlabel(xlabel)
    ax.set_ylabel(ylabel)
    ax.set_xlim(0, xlimit)
    fig = ax.get_figure()
    fig.set_size_inches(10, 7)
    if path:
         plt.savefig(path + '//' + save_as, bbox_inches = 'tight', dpi = 600)
    plt.show()
    plt.close()

# Cell
@exception_handler
def pointplot_plate(df = None, x = None, y = None, hue = None, hue_order = None, threshold = None, ylabel = None,
                    palette = None,  save_as = None, path  = None):
    """Creates point plot for the data.
    :param df: pandas DataFrame for plotting.
    :param x, y: names of variables in data.
    :param hue, hue_order, ylabel, xlimit, palette: These features are seaborn.boxplot features.
    :param threshold: threshold for hit identification.
    :param save_as: filename to save the resulting figure, optional.
    :param path: path to the output folder to save the resulting figure, optional.
    :return: None."""
    import pandas as pd
    import matplotlib as mpl
    import matplotlib.colors as colors
    import matplotlib.pyplot as plt
    import seaborn as sns

    g = sns.catplot(data = df, x = x, y = y, hue = hue, height = 6, aspect = 2.5, margin_titles = False,
                   palette = palette, hue_order = hue_order)
    if threshold:
        plt.plot([0, len(df.Well.unique())], [threshold, threshold],'r-')
        plt.plot([0, len(df.Well.unique())], [-threshold, -threshold],'r-')

    g.set_xticklabels([])
    plt.ylabel(ylabel)
    g.despine()
    if path:
         g.savefig(path +'//' + save_as, bbox_inches = 'tight', dpi = 600)
    plt.show()
    plt.close()

# Cell
@exception_handler
def calculate_growth_score(df = None):
    """Calculates growth scores from time series data.
    :param df: pandas DataFrame with time-series data.
    :return: DataFrame with growth scores."""
    import pandas as pd
    import numpy as np

    df = df.astype(float).sort_values(['Time']) # sort by time
    times = df['Time'].values.astype(int)
    # create time-series table
    ts_data = pd.DataFrame()
    for name, data in df.drop(columns = ['Time']).iteritems():
        ts_data = ts_data.append(pd.DataFrame({'Well': name, 'Time': times, 'OD': data}))

    # calculate drowth score
    score_data = pd.DataFrame()
    for name, well in ts_data.groupby('Well'):
        well = well.copy()
        well['past'] = np.append(well['OD'].values[0], well['OD'].values[:-1])
        well['grate'] = (well['OD'] - well['past'])/well['past']
        well['gscore'] = (well['OD'].max() - well['OD'].values[0]) + well['grate'].max()*0.25
        score_data = score_data.append(well)

    score_data = score_data.drop(columns = ['past'])
    return score_data

# Cell
@exception_handler
def filter_curves(df = None):
    """Filter out aberrant curves.
    :param df: pandas DataFrame with time-series data.
    :return: clean DataFrame with growth scores."""
    import pandas as pd

    clean = pd.DataFrame()
    for name, well in df.groupby('Well'):
        well = well.copy().sort_values(['Time']).reset_index()

        # if there is big sudden drop
        if (well['grate'].min() < -0.2) and (well['grate'].idxmin() > 4): #0.1
            print(name, well['grate'].min())
            well['Result'] = 'Invalid_sample'

        # if the curve started at high value
        elif well['OD'].values[0] >  0.2:
            print(name)
            well['Result'] = 'Invalid_sample'
        else:
            well['Result'] = well['Status']

        clean = clean.append(well)
    return clean

# Cell
@exception_handler
def inv_log(x = None):
    """Inverse log calculator"""
    return ((10**-x)/(1e-6))

# Cell
@exception_handler
def ll4(x,b,c,d,e):
    """Fitting function - LM equation, LL.4 function (4-parameter sigmoidal function).
     - b: hill slope
     - c: min response
     - d: max response
     - e: EC50"""
    import numpy as np
    import warnings
    warnings.filterwarnings('ignore')
    return(c+(d-c)/(1+np.exp(b*(np.log(x)-np.log(e)))))

# Cell
@exception_handler
def pDose(x = None):
    """Helper function, compute log transformed concentrations."""
    import numpy as np
    return(-np.log10(1e-6*x))

# Cell
@exception_handler
def dose_response(df = None, y_label = 'Response', path = None):
    """Dose response function. The input DataFrame should contain columns 'Compound_id', 'Dose', 'Response'.
    The DataFrame shouldn't contain NAN values or dose 0, which will result in infinity at logDose.
    The fitting function is a LL.4 function (4-parameter sigmoidal function) with
     - b: hill slope
     - c: min response
     - d: max response
     - e: EC50
    :param df: pandas DataFrame for plotting.
    :param y_label: name for y-axis.
    :param path: path to the output folder to save the results, optional.
    :return: DataFrame with Dose Response fitting results."""
    import pandas as pd
    import numpy as np
    import scipy.optimize as opt
    from scipy.stats.stats import pearsonr
    import matplotlib as mpl
    import matplotlib.colors as colors
    import matplotlib.pyplot as plt
    import seaborn as sns
    sns.set(context = 'notebook', style = 'white', palette = 'dark')
    import warnings
    warnings.filterwarnings('ignore')

    pDose = lambda x:-np.log10(1e-6*x)
    ll4 = lambda x, b, c, d, e:(c+(d-c)/(1+np.exp(b*(np.log(x)-np.log(e)))))
    df = df[['Compound_id', 'Dose', 'Response']].copy()
    df = df[(df != 0).all(1)]  # drop zero values
    df['logDose'] = pDose(df.Dose.astype(float)) # calculate logDose
    df_mean = df.groupby(['Compound_id','Dose'], as_index = False).mean() # calculate response mean values
    df_mean['std'] = list(df.groupby(['Compound_id','Dose']).std()['Response'].values) # calculate response std

    fitData = pd.DataFrame()
    for name, group in df_mean.groupby(['Compound_id']): # group data by compounds
        print(name)
    # fitting curve
        try:
                fitCoefs, covMatrix = opt.curve_fit(ll4, group.Dose, group.Response, method = 'lm')
                residuals = group.Response - group.Dose.apply(lambda x: ll4(x,*fitCoefs))
                curFit = dict(zip(['b','c','d','e'], fitCoefs))
                curFit['Compound_id'], curFit['residuals'] = name, sum(residuals**2)
                predicted = group.Dose.apply(lambda x: ll4(x,*fitCoefs))
                curFit['r_squared'] = pearsonr(group.Response, predicted)[0]**2
                curFit['N'] =  int(group.shape[0])
                fitData = fitData.append(curFit, ignore_index = True)
                EC50_response = ll4(curFit['e'],*[curFit[i] for i in ['b','c','d','e']])

                # plot data
                raw = df[df.Compound_id == name]
                refDose = np.linspace(min(raw.Dose)*0.55, max(raw.Dose)*1.6, 256)
                g2 = sns.lmplot('logDose', 'Response', data = group,  fit_reg = False, legend = False, height=6)
                g2.map(plt.errorbar, 'logDose', 'Response',yerr = group['std'], fmt='o')
                axes = plt.gca()
                axes.invert_xaxis()
                plt.plot([pDose(i) for i in refDose],[ll4(i,*[curFit[i] for i in ['b','c','d','e']]) for i in refDose])
                locs, labels = plt.xticks()
                g2.set_xticklabels([round(inv_log(l), 1) for l in locs]) # inverse log for xticks
                plt.xlabel('Dose (um)')
                plt.ylabel(y_label)
                plt.title(name)

                #plot EC_50_label
                ymin, ymax = axes.get_ylim()
                xmin, xmax = axes.get_xlim()
                plt.plot([xmin, pDose(curFit['e'])], [EC50_response, EC50_response], color = 'navy', linestyle = '--', lw = 0.7)
                plt.plot([pDose(curFit['e']), pDose(curFit['e'])], [ymin, EC50_response], color = 'navy', linestyle = '--', lw = 0.7)
                plt.show()

                if path:
                    g2.savefig(path +'//' + name +'_dr_fit.png', bbox_inches='tight', dpi=600)

                plt.close()

        except Exception as e:
                print('Fitting curve failed:')
                print(e)
    if not fitData.empty:
        fitData = fitData.set_index('Compound_id') # .round(2)
        fitData['N'] = fitData.N.astype(int)
        fitData.rename(columns = {'b': 'hill slope', 'c': 'min response', 'd': 'max response', 'e': 'EC50'}, inplace = True)
        return fitData

# Cell
@exception_handler
def dose_response_old(df = None, y_label = 'Response', path = None):
    """Dose response function. The input DataFrame should contain columns 'Compound_id', 'Dose', 'Response'.
    The DataFrame shouldn't contain NAN values or dose 0, which will result in infinity at logDose.
    :param df: pandas DataFrame for plotting.
    :param y_label: name for y-axis.
    :param path: path to the output folder to save the results, optional.
    :return: DataFrame with Dose Response fitting results."""
    import pandas as pd
    import numpy as np
    import scipy.optimize as opt
    from scipy.stats.stats import pearsonr
    import matplotlib as mpl
    import matplotlib.colors as colors
    import matplotlib.pyplot as plt
    import seaborn as sns
    sns.set(context = 'notebook', style = 'white', palette = 'dark')
    import warnings
    warnings.filterwarnings('ignore')

    pDose = lambda x:-np.log10(1e-6*x)
    ll4 = lambda x, b, c, d, e:(c+(d-c)/(1+np.exp(b*(np.log(x)-np.log(e)))))


    df = df[['Compound_id', 'Dose', 'Response']].copy()
    df = df[(df != 0).all(1)]  # drop zero values
    df['logDose'] = pDose(df.Dose.astype(float)) # calculate logDose
    df_mean = df.groupby(['Compound_id','Dose'], as_index = False).mean() # calculate response mean values
    df_mean['std'] = list(df.groupby(['Compound_id','Dose']).std()['Response'].values) # calculate response std

    fitData = pd.DataFrame()
    for name, group in df_mean.groupby(['Compound_id']): # group data by compounds
        print(name)
    # fitting curve
        try:
                fitCoefs, covMatrix = opt.curve_fit(ll4, group.Dose, group.Response, method = 'lm')
                residuals = group.Response - group.Dose.apply(lambda x: ll4(x,*fitCoefs))
                curFit = dict(zip(['b','c','d','e'], fitCoefs))
                curFit['Compound_id'], curFit['residuals'] = name, sum(residuals**2)
                predicted = group.Dose.apply(lambda x: ll4(x,*fitCoefs))
                curFit['r_squared'] = pearsonr(group.Response, predicted)[0]**2
                curFit['N'] =  int(group.shape[0])
                fitData = fitData.append(curFit, ignore_index = True)
                EC50_response = ll4(curFit['e'],*[curFit[i] for i in ['b','c','d','e']])


                # plot data
                raw = df[df.Compound_id == name]
                refDose = np.linspace(min(raw.Dose)*0.55, max(raw.Dose)*1.6, 256)
                g2 = sns.lmplot('logDose', 'Response', data = group,  fit_reg = False, legend = False, height=6)
                g2.map(plt.errorbar, 'logDose', 'Response',yerr = group['std'], fmt='o')
                #plt.ylim(group['Response'].min()*1.5, group['Response'].max()*1.2)
                plt.xlim(max(raw.logDose)*1.1, min(raw.logDose)*0.9)
                plt.plot([pDose(i) for i in refDose],[ll4(i,*[curFit[i] for i in ['b','c','d','e']]) for i in refDose])
                locs, labels = plt.xticks()
                g2.set_xticklabels([round(inv_log(l), 1) for l in locs]) # inverse log for xticks
                plt.xlabel('Dose (um)')
                plt.ylabel(y_label)
                plt.title(name)
                #plot EC_50_label
                #plt.plot([max(raw.logDose)*1.1, pDose(curFit['e'])], [EC50_response, EC50_response], color = 'navy', linestyle = '--', lw = 0.7)
                #plt.plot([pDose(curFit['e']), pDose(curFit['e'])], [group['Response'].min()*1.2, EC50_response], color = 'navy', linestyle = '--', lw = 0.7)
                plt.plot([max(raw.logDose)*1.1, pDose(curFit['e'])], [EC50_response, EC50_response], color = 'navy', linestyle = '--', lw = 0.7)
                plt.plot([pDose(curFit['e']), pDose(curFit['e'])], [group['Response'].min()*0.6, EC50_response], color = 'navy', linestyle = '--', lw = 0.7)

                plt.show()
                if path:
                    g2.savefig(path +'//' + name +'_dr_fit.png', bbox_inches='tight', dpi=600)

                plt.close()

        except Exception as e:
                print('Fitting curve failed:')
                print(e)
    if not fitData.empty:
        fitData = fitData.set_index('Compound_id') # .round(2)
        fitData['N'] = fitData.N.astype(int)
        fitData.rename(columns = {'b': 'hill slope', 'c': 'min response', 'd': 'max response', 'e': 'EC50'}, inplace = True)
        return fitData

# Cell
@exception_handler
def plot_dr_viability(data = None, y_label = 'Response', path = None):
    """Plots response vs viability. The DataFrame should contain columns ['Compound', 'Dose','logDose', 'Viability', 'Response'] (at least)."""
    import pandas as pd
    import matplotlib as mpl
    import matplotlib.colors as colors
    import matplotlib.pyplot as plt
    import seaborn as sns
    sns.set(context = 'notebook', style = 'white', palette = 'dark')

    df = data[['Compound_id', 'Dose','logDose', 'Viability', 'Response']]
    df = df[(df != 0).all(1)]  # drop zero values
    df_mean = df.groupby(['Compound_id','Dose'], as_index = False).mean() # calculate response mean values
    df_mean['resp_std'] = list(df.groupby(['Compound_id','Dose']).std()['Response'].values) # calculate response std
    df_mean['via_std'] = list(df.groupby(['Compound_id','Dose']).std()['Viability'].values) # calculate viability std

    for name, group in  df_mean.groupby('Compound_id'):  # group data by compounds
        group = group.sort_values('Dose')
        error_resp, error_via  = group['resp_std'], group['via_std']

        fig, ax1 = plt.subplots(figsize = (6,6))
        plt.title(name, fontsize = 16)

        plot1 = ax1.plot(group['logDose'], group['Response'], 'b', label = 'Response')
        ax1.set_xlim(max(group['logDose'])*1.07, min(group['logDose'])*0.9)
        ax1.set_ylabel(y_label, fontsize = 16)
        ax1.set_ylim(0, df_mean['Response'].max()*1.2)
        ax1.errorbar(group['logDose'], group['Response'],yerr = error_resp, fmt ='o', color ='b', ecolor = 'lightblue')

        ax2 = ax1.twinx()
        plot2 = ax2.plot(group['logDose'], group['Viability'], 'g', label = 'Viability')
        ax2.set_xlim(max(group['logDose'])*1.07, min(group['logDose'])*0.9)
        ax2.set_ylabel('Viability', fontsize = 16)
        ax2.set_ylim(0, 120)
        ax2.errorbar(group['logDose'], group['Viability'],yerr=error_via, fmt='o', color='g', ecolor='lightgreen')
        ax1.set_xlabel('Dose, um', fontsize = 16)

        # create legend
        lines = plot1 + plot2
        ax1.legend(lines, [l.get_label() for l in lines])
        locs, labels = plt.xticks()
        new_labels =[]
        for loc in locs:
            inv_log = lambda x:((10**-x)/(1e-6)) # inverse log calculator to set xticks
            inv_x = round(inv_log(loc), 1)
            new_labels.append(inv_x)
        ax1.set_xticklabels(new_labels)
        if path:
            plt.savefig(path +'//' + name +'_raw_viability.png', bbox_inches='tight', dpi=600)
        plt.show()

# Cell
@exception_handler
def prune_compound(df, threshold = -0.15):
    """This function takes DataFrame of one-compound dose-response data, find maximum activity,
    and drops rows starting from treshold-defined reduction of Response. The default value for threshold = -0.15,
   it drops rows starting from 15% reduction of Response. The input DataFrame should contain columns
   'Compound_id', 'Dose', 'Response'."""

    prunned = pd.DataFrame()
    df = df.sort_values('Dose')
    curr_max = 0.0000001
    groups = df.groupby('Dose')
    for name, group in groups:
        percent_change = (group['Response'].mean()/curr_max)-1
        if group['Response'].mean() > curr_max:
            curr_max = group['Response'].mean()
        if percent_change > threshold:
            prunned = prunned.append(group)
    return(prunned)

# Cell
@exception_handler
def plot_polynomial(df = None, y_label = 'Response', degree = 2, path = None):
    """Plot polynomial fit.
    :param df: pandas DataFrame for plotting.
    :param y_label: name for y-axis.
    :param degree: degree of the polynomial fit.
    :param path: path to the output folder to save the results, optional.
    :return: None."""
    import pandas as pd
    import numpy as np
    import matplotlib as mpl
    import matplotlib.colors as colors
    import matplotlib.pyplot as plt
    import seaborn as sns

    try:
        df = df[['Compound_id', 'Dose', 'Response']]
        df = df[(df != 0).all(1)]  # drop zero values
        df['logDose'] = pDose(df['Dose'].astype(float))
    except Exception as e:
            print(e)

    # plot response
    for name, group in df.groupby('Compound_id'):
        mean_group = group.groupby(['Dose'], as_index = False).mean()
        g = sns.lmplot('logDose', 'Response', data = mean_group,  fit_reg = False, legend = False, height=6)
        g.map(plt.errorbar, 'logDose', 'Response',yerr = list(group.groupby(['Dose']).std()['Response'].values), fmt='o')
        plt.xlim(max(group['logDose'])*1.1, min(group['logDose'])*0.9)
        locs, labels = plt.xticks()
        g.set_xticklabels([round(inv_log(l), 1) for l in locs]) # inverse log for xticks
        plt.xlabel('Dose (um)')
        plt.ylabel(y_label)
        plt.title(name)

        # plot polynomial_fit
        try:
            X  = np.asarray(group['logDose'].values)
            Y  = np.asarray(group['Response'].values)
            p = np.poly1d(np.polyfit(X,Y,degree))
            polyDose = np.linspace(min(group['logDose'])*0.98, max(group['logDose'])*1.02, 256)
            plt.plot(polyDose, p(polyDose), color='Navy')

        except Exception as e:
            print('Polynomial fit failed:')
            print(e)
        plt.show()
        if path:
            g.savefig(path +'//' + name +'_polynomial.png', bbox_inches='tight', dpi=600)
        plt.close()

# Cell
@exception_handler
def df_to_table(df = None, slide = None, left = None, top = None, width = None, height = None, colnames = None):
    """Converts a Pandas DataFrame to a PowerPoint table on the given slide of a PowerPoint presentation.
    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools (resizing columns, changing formatting etc).
    Source:  https://github.com/robintw/PandasToPowerpoint/blob/master/PandasToPowerpoint.py
    :param df: pandas DataFrame with the data.
    :param slide: slide object from the python-pptx library containing the slide on which you want the table to appear
    :param left, top, right, width, height, colnames: These parameters are python-pptx parameters.
    :return: Powerpoint table.
     """
    import pandas as pd
    from pptx import Presentation
    from pptx.util import Inches, Pt

    rows, cols = df.shape
    res = slide.shapes.add_table(rows + 1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)
        colnames[0] = 'idx'

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        # Column names can be tuples
        if not isinstance(col_name, str):
            col_name = ' '.join(col_name)
        res.table.cell(0, col_index).text = col_name
        paragraph = res.table.cell(0, col_index).text_frame.paragraphs[0]
        paragraph.font.size = Pt(9)

    #m = df.as_matrix()
    m = df.values

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            res.table.cell(row + 1, col).text = text
            paragraph = res.table.cell(row + 1, col).text_frame.paragraphs[0]
            paragraph.font.size = Pt(9)

# Cell
@exception_handler
def create_presentation(path = None):
    """Creates ppt report from files in the specified folder.
    Reads template ppt file and files in the input folder, adds all pictures and tables with
    size less than 30 rows from input folder.
    :param path: path to the output folder to save the results.
    :return: None.
    """
    import os
    import pandas as pd
    from datetime import date
    from pptx import Presentation
    from pptx.util import Inches, Pt

    report = Presentation('hts_data//templates//ppt_template.pptx')
    slide = report.slides.add_slide(report.slide_layouts[6])
    #pic = slide.shapes.add_picture('hts_data//templates//logo.png', left = Inches(3), top = Inches(0.2))
    subtitle = slide.shapes.add_textbox(left = Inches(5.), top = Inches(3.5), width = Inches(3), height = Inches(0.5),).text_frame
    p = subtitle.paragraphs[0]
    run = p.add_run()
    run.text = 'Technical Report\nGenerated on {:%m-%d-%Y}'.format(date.today())
    font = run.font
    font.size = Pt(18)
    files_list = os.listdir(path)
    for myfile in files_list:
        if 'heatmap.png' in myfile:
            slide = report.slides.add_slide(report.slide_layouts[6])
            left = top = Inches(0.7)
            height = Inches(6)
            pic = slide.shapes.add_picture(path + '//' + myfile, left, top, width = Inches(5.8), height= Inches(4))
        elif '.png' in myfile and 'heatmap.png' not in myfile:
            slide = report.slides.add_slide(report.slide_layouts[6])
            subtitle = slide.shapes.add_textbox(left = Inches(0.5), top = Inches(0.3), width = Inches(2), height = Inches(0.5)).text_frame
            subtitle.text = myfile
            left = top = Inches(0.7)
            pic = slide.shapes.add_picture(path +'//' + myfile, left, top = Inches(0.8), height= Inches(6))
            left = Inches(0.7)
        elif 'csv' in myfile:
            try:
                table = pd.read_csv(path +'//' + myfile)
                if table.shape[0]<30:
                    slide = report.slides.add_slide(report.slide_layouts[6])
                    subtitle = slide.shapes.add_textbox(left = Inches(0.5), top = Inches(0.3), width = Inches(2), height = Inches(0.5)).text_frame
                    subtitle.text = myfile
                    slide_table = df_to_table(table, slide, left = Inches(0.3), top = Inches(1), width = Inches(12.5), height = Inches(0.3))
                left = Inches(0.7)
            except Exception as e:
                print(e)
    return report
